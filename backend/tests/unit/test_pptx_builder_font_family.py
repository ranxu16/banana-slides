from pptx.oxml.ns import qn

from services.image_editability.text_attribute_extractors import ColoredSegment, TextStyleResult
from utils.pptx_builder import PPTXBuilder


def _build_slide():
    builder = PPTXBuilder()
    presentation = builder.create_presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    return builder, slide


def test_single_color_text_sets_east_asian_font_family():
    builder, slide = _build_slide()

    builder.add_text_element(
        slide=slide,
        text='中文字体',
        bbox=[0, 0, 320, 80],
        text_style=TextStyleResult(
            font_family='Microsoft YaHei',
            font_color_rgb=(17, 34, 51),
        ),
    )

    run = slide.shapes[-1].text_frame.paragraphs[0].runs[0]
    run_properties = run._r.get_or_add_rPr()

    assert run.font.name == 'Microsoft YaHei'
    assert run_properties.find(qn('a:latin')).get('typeface') == 'Microsoft YaHei'
    assert run_properties.find(qn('a:ea')).get('typeface') == 'Microsoft YaHei'


def test_multi_color_text_sets_east_asian_font_family_on_each_run():
    builder, slide = _build_slide()

    builder.add_text_element(
        slide=slide,
        text='ignored',
        bbox=[0, 0, 320, 80],
        text_style=TextStyleResult(
            font_family='Microsoft YaHei',
            colored_segments=[
                ColoredSegment(text='中', color_rgb=(255, 0, 0)),
                ColoredSegment(text='文', color_rgb=(0, 0, 255)),
            ],
        ),
    )

    runs = slide.shapes[-1].text_frame.paragraphs[0].runs
    assert len(runs) == 2

    for run in runs:
        run_properties = run._r.get_or_add_rPr()
        assert run.font.name == 'Microsoft YaHei'
        assert run_properties.find(qn('a:latin')).get('typeface') == 'Microsoft YaHei'
        assert run_properties.find(qn('a:ea')).get('typeface') == 'Microsoft YaHei'


def test_font_family_nodes_are_inserted_before_hyperlink_properties():
    builder, slide = _build_slide()

    builder.add_text_element(
        slide=slide,
        text='中文字体',
        bbox=[0, 0, 320, 80],
        text_style=TextStyleResult(font_family='Microsoft YaHei'),
    )

    run = slide.shapes[-1].text_frame.paragraphs[0].runs[0]
    run_properties = run._r.get_or_add_rPr()
    run_properties.get_or_add_hlinkClick()
    PPTXBuilder._apply_font_family_to_run(run, 'Microsoft YaHei')

    child_tags = [child.tag for child in run_properties]
    assert child_tags.index(qn('a:latin')) < child_tags.index(qn('a:hlinkClick'))
    assert child_tags.index(qn('a:ea')) < child_tags.index(qn('a:hlinkClick'))
